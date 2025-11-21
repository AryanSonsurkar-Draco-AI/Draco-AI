# main.py
# Draco: unified Flask + SocketIO assistant backend
# Save next to draco.html and run: python main.py
# Required packages (example):
# pip install flask flask-socketio pyttsx3 pygame SpeechRecognition ddgs psutil pyautogui pywhatkit
import os
import sys
import time
import json
import random
import datetime
import threading
import subprocess
import platform
import webbrowser
from urllib.parse import quote as url_quote
from urllib.parse import urlparse
from collections import deque
from typing import Optional
import re
from flask import session
from flask import Flask, send_from_directory, request, session, redirect, url_for, jsonify
from flask_socketio import SocketIO, emit
from werkzeug.utils import secure_filename
try:
    from docx import Document
except Exception:
    Document = None
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
except Exception:
    Presentation = None
    Inches = None
    Pt = None
    PP_ALIGN = None
    RGBColor = None
    MSO_SHAPE = None
try:
    from fpdf import FPDF
except Exception:
    FPDF = None
try:
    from PyPDF2 import PdfReader
except Exception:
    PdfReader = None
import sympy as sp
# optional imports (best-effort)
try:
    import pyttsx3
except Exception:
    pyttsx3 = None
try:
    import pygame
except Exception:
    pygame = None
try:
    import speech_recognition as sr
except Exception:
    sr = None
if sr:
    try:
        r = sr.Recognizer()
    except Exception:
        r = None
try:
    import psutil
except Exception:
    psutil = None
try:
    import pyautogui
except Exception:
    pyautogui = None
try:
    import pywhatkit
except Exception:
    pywhatkit = None
# DuckDuckGo scraping wrapper
try:
    from ddgs import DDGS
except Exception:
    DDGS = None
# Music library hook (user's file)
try:
    import musicLibrary
except Exception:
    musicLibrary = None
import requests
import pytz
try:
    import draco_chat
except Exception:
    draco_chat = None

try:
    from draco_intro import handle_intro as _handle_intro_request
except Exception:
    _handle_intro_request = None

# Disable legacy draco_chat fallback responses (keep routing in main.py only)
draco_chat = None

BossLevelEngine = None  # Boss level module disabled

ON_RENDER = os.environ.get("RENDER") is not None
ON_SERVER = ON_RENDER or (os.environ.get("PORT") is not None) or (os.environ.get("RENDER_EXTERNAL_URL") is not None)

# ------------- Flask / SocketIO -------------
# Force Flask-SocketIO to use threading instead of eventlet or gevent
os.environ["FLASK_SOCKETIO_ASYNC_MODE"] = "threading"
app = Flask(__name__, static_folder=".", template_folder=".")
app.secret_key = os.environ.get("FLASK_SECRET_KEY", "dev-secret-change-me")
socketio = SocketIO(app, cors_allowed_origins="*", async_mode="threading")
# Limit uploads to 20 MB
app.config["MAX_CONTENT_LENGTH"] = 20 * 1024 * 1024

@app.route("/")
def home():
    """Serve the main Draco web client."""
    return send_from_directory(".", "draco.html")

# ------------- Global utilities & config -------------
MEMORY_FILE = "memory.json"
NOTES_FILE = "notes.json"
REMINDERS_FILE = "reminders.json"
WEATHER_API_KEY = ""   # add your OpenWeatherMap key if desired
NEWSAPI_KEY = ""       # add your News API key if desired
USERS_DIR = os.path.join(os.getcwd(), "users")

def safe_write_json(path, data):
    with open(path, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2)

def safe_read_json(path, default):
    if os.path.exists(path):
        try:
            with open(path, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception:
            return default
    return default

def ensure_dir(path):
    try:
        os.makedirs(path, exist_ok=True)
    except Exception:
        pass

def sanitize_email_for_path(email: str) -> str:
    # basic safe filename for email
    return email.replace("@", "_at_").replace(".", "_dot_")

def user_paths(email: str):
    root = os.path.join(USERS_DIR, sanitize_email_for_path(email))
    ensure_dir(root)
    return {
        "root": root,
        "profile": os.path.join(root, "profile.json"),
        "chat": os.path.join(root, "chat.json"),  # legacy single-thread file
        "chats": os.path.join(root, "chats.json"), # new: multi-chat storage
    }

def get_logged_in_email() -> Optional[str]:
    e = session.get("user_email")
    if isinstance(e, str) and "@" in e:
        return e
    return None

def _load_chats(email: str):
    p = user_paths(email)
    chats = safe_read_json(p["chats"], None)
    if chats is None:
        # migrate from legacy chat.json into a single chat
        legacy = safe_read_json(p["chat"], [])
        chat_id = str(int(time.time() * 1000))
        name = _summarize_chat_name(legacy)
        chats = [{"id": chat_id, "name": name, "items": legacy, "updated_at": time.time()}]
        safe_write_json(p["chats"], chats)
    return chats

def _save_chats(email: str, chats):
    p = user_paths(email)
    safe_write_json(p["chats"], chats)

def _current_chat_id() -> Optional[str]:
    cid = session.get("chat_id")
    return str(cid) if cid else None

def _set_current_chat_id(cid: str):
    session["chat_id"] = str(cid)

def _get_or_create_current_chat(email: str):
    chats = _load_chats(email)
    cid = _current_chat_id()
    chat = None
    if cid:
        for c in chats:
            if c.get("id") == cid:
                chat = c
                break
    if chat is None:
        # pick latest or create
        if chats:
            chats = sorted(chats, key=lambda c: c.get("updated_at", 0), reverse=True)
            chat = chats[0]
            _set_current_chat_id(chat.get("id"))
        else:
            chat = _create_new_chat(email)
            chats = _load_chats(email)  # ensure persisted
    return chat

def _create_new_chat(email: str):
    chats = _load_chats(email)
    cid = str(int(time.time() * 1000))
    chat = {"id": cid, "name": "New chat", "items": [], "updated_at": time.time()}
    chats.append(chat)
    _save_chats(email, chats)
    _set_current_chat_id(cid)
    return chat

def _clear_current_chat(email: str):
    chats = _load_chats(email)
    cid = _current_chat_id()
    for c in chats:
        if c.get("id") == cid:
            c["items"] = []
            c["updated_at"] = time.time()
            _save_chats(email, chats)
            return True
    return False

STOPWORDS = set("the a an and or to for in of on with at by from as be is are was were it this that these those i you he she we they them our your my mine ours yours their his her its about into than too very just can could should would will shall may might do does did not no yes ok please hey hi hello how what when where which who why".split())

def _summarize_chat_name(items):
    # items: list of {who, text}
    texts = [x.get("text", "") for x in items if x.get("who") == "user"]
    if not texts:
        texts = [x.get("text", "") for x in items]
    combined = " ".join(texts).strip()
    if not combined:
        return datetime.datetime.now().strftime("Chat %b %d %H:%M")
    # take first sentence
    first = combined.split("\n")[0].split(".")[0]
    words = [w.strip(" ,:;!?()[]{}\"'\t").lower() for w in first.split()]
    filtered = [w for w in words if w and w not in STOPWORDS and all(ch.isalnum() or ch in "-_'" for ch in w)]
    if not filtered:
        # fallback: first few original words
        filtered = first.split()
    # min 2 words, max 10 words
    if len(filtered) < 2:
        extra = (combined.split())[0:2]
        filtered = (filtered + extra)[:2]
    name_words = filtered[:10]
    name = " ".join(name_words).strip().title()
    if len(name.split()) < 2:
        name = (name + " Chat").strip()
    return name[:80]

def save_chat_line(email: str, who: str, text: str):
    chats = _load_chats(email)
    chat = _get_or_create_current_chat(email)
    # refresh chats after current selection (IDs)
    cid = chat.get("id")
    for c in chats:
        if c.get("id") == cid:
            c.setdefault("items", []).append({"ts": time.time(), "who": who, "text": text})
            # set name on first meaningful lines
            if (not c.get("name")) or c.get("name") in ("New chat", ""):
                c["name"] = _summarize_chat_name(c["items"]) or "New Chat"
            c["updated_at"] = time.time()
            break
    _save_chats(email, chats)

def get_chat_history(email: str, chat_id: Optional[str] = None):
    chats = _load_chats(email)
    cid = chat_id or _current_chat_id()
    if cid:
        for c in chats:
            if c.get("id") == cid:
                return c.get("items", [])
    # default to latest
    if chats:
        chats = sorted(chats, key=lambda c: c.get("updated_at", 0), reverse=True)
        _set_current_chat_id(chats[0].get("id"))
        return chats[0].get("items", [])
    return []

def get_user_profile(email: str):
    p = user_paths(email)
    return safe_read_json(p["profile"], {})

def set_user_profile(email: str, profile: dict):
    p = user_paths(email)
    safe_write_json(p["profile"], profile or {})

# ------------- Memory / Personality -------------
class MemoryManager:
    def __init__(self, memory_file=MEMORY_FILE):
        self.memory_file = memory_file
        self.session = deque(maxlen=60)
        self.long = safe_read_json(memory_file, {"language": "english", "name": "Ars", "mood": "neutral"})
        self._save()

    def add(self, text):
        self.session.append({"time": time.time(), "text": text})
        self._save()

    def get_session(self):
        return list(self.session)

    def set_pref(self, k, v):
        self.long[k] = v
        self._save()

    def get_pref(self, k, default=None):
        return self.long.get(k, default)

    def _save(self):
        safe_write_json(self.memory_file, self.long)

class Personality:
    def __init__(self):
        self.emotion = "neutral"

    def update(self, text: str):
        t = text.lower()
        if any(x in t for x in ["happy", "great", "awesome", "nice"]):
            self.emotion = "happy"
        elif any(x in t for x in ["sad", "upset", "bad", "angry"]):
            self.emotion = "concerned"
        else:
            # small random drift
            self.emotion = random.choice(["neutral", "playful", "helpful"])

    def respond(self, base: str) -> str:
        # dynamic replies based on emotion
        if self.emotion == "happy":
            templates = [f"Yay! {base}", f"I'm glad — {base}", f"Nice! {base}"]
        elif self.emotion == "concerned":
            templates = [f"I hear you. {base}", f"I'm here for you. {base}", f"Don't worry — {base}"]
        elif self.emotion == "playful":
            templates = [f"Heh — {base}", f"Alrighty! {base}", f"Let's do it: {base}"]
        else:
            templates = [base, f"Okay. {base}", f"Done — {base}"]
        return random.choice(templates)

memory = MemoryManager()
personality = Personality()
chat_ctx = None

# ------------- Gamification (XP / Levels / Modes) -------------

def _ensure_gamification():
    g = memory.long.get("gamification")
    if not isinstance(g, dict):
        g = {}
    if "level" not in g:
        g["level"] = 1
    if "xp" not in g:
        g["xp"] = 0
    if "mode" not in g:
        g["mode"] = "balanced"  # study / coding / balanced
    if "total_xp" not in g:
        g["total_xp"] = 0
    if "last_action" not in g:
        g["last_action"] = None
    if "badges" not in g or not isinstance(g["badges"], list):
        g["badges"] = []
    memory.long["gamification"] = g
    memory._save()
    return g


def _xp_needed(level: int) -> int:
    try:
        lvl = max(1, int(level))
    except Exception:
        lvl = 1
    return 50 + lvl * 20


def add_xp(kind: str, amount: int):
    """Award XP for an action kind (study/coding/general)."""
    g = _ensure_gamification()

    if amount <= 0:
        return {
            "gained": 0,
            "levels_gained": 0,
            "level": g.get("level", 1),
            "xp": g.get("xp", 0),
            "xp_needed": _xp_needed(g.get("level", 1)),
            "total_xp": g.get("total_xp", 0),
            "kind": kind,
        }

    # small boost based on current mode
    mode = g.get("mode", "balanced")
    boost = 1.0

    if mode == "study" and kind == "study":
        boost = 1.2
    elif mode == "coding" and kind == "coding":
        boost = 1.2

    gained = int(amount * boost)

    g["xp"] = int(g.get("xp", 0)) + gained
    g["total_xp"] = int(g.get("total_xp", 0)) + gained
    g["last_action"] = {"kind": kind, "amount": gained, "ts": time.time()}

    # Level up loop
    levels_gained = 0
    while g["xp"] >= _xp_needed(g["level"]):
        needed = _xp_needed(g["level"])
        g["xp"] -= needed
        g["level"] += 1
        levels_gained += 1

    memory.long["gamification"] = g
    memory._save()

    result = {
        "gained": gained,
        "levels_gained": levels_gained,
        "level": g["level"],
        "xp": g["xp"],
        "xp_needed": _xp_needed(g["level"]),
        "total_xp": g.get("total_xp", 0),
        "kind": kind,
    }

    if levels_gained:
        msg = f" Level up! You're now Level {g['level']} (XP {g['xp']}/{result['xp_needed']})."
        memory.add(f"System: {msg}")
        emit_to_ui("draco_response", {"text": msg})
        emit_to_ui("level_up", {
            "level": g["level"],
            "xp": g["xp"],
            "xp_needed": result["xp_needed"],
            "kind": kind,
        })

    email = get_logged_in_email()
    if email:
        try:
            prof = get_user_profile(email)
            prof["level"] = g["level"]
            prof["xp"] = g.get("total_xp", 0)
            set_user_profile(email, prof)
        except Exception:
            pass

    return result


XP_KEYWORD_RULES = [
    {"kind": "study", "amount": 6, "terms": ["study", "exam", "test", "revision", "revise", "assignment", "homework", "flashcard", "notes", "note ", "note", "memorize", "syllabus"]},
    {"kind": "study", "amount": 5, "terms": ["research", "summarize", "summary", "analysis", "analyze", "outline", "explain", "reference"]},
    {"kind": "coding", "amount": 6, "terms": ["code", "coding", "program", "script", "develop", "feature", "algorithm", "function"]},
    {"kind": "coding", "amount": 7, "terms": ["bug", "error", "stack trace", "issue", "crash", "fix", "broken", "failure", "exception"]},
    {"kind": "coding", "amount": 8, "terms": ["refactor", "cleanup", "optimize", "improve structure", "rewrite", "tidy code"]},
    {"kind": "study", "amount": 4, "terms": ["challenge", "quest", "focus", "grind", "level up", "xp"]},
]


def _maybe_award_keyword_xp(cmd: str):
    if not cmd:
        return
    lower = cmd.lower()
    for rule in XP_KEYWORD_RULES:
        if any(term in lower for term in rule["terms"]):
            add_xp(rule["kind"], rule["amount"])


def get_gamestate():
    g = _ensure_gamification()
    lvl = int(g.get("level", 1))
    xp = int(g.get("xp", 0))
    needed = _xp_needed(lvl)
    return {
        "level": lvl,
        "xp": xp,
        "xp_needed": needed,
        "mode": g.get("mode", "balanced"),
        "total_xp": int(g.get("total_xp", 0)),
        "badges": g.get("badges", []),
        "last_action": g.get("last_action"),
    }


def set_mode(mode: str):
    g = _ensure_gamification()
    m = (mode or "balanced").lower()
    if m not in ("study", "coding", "balanced"):
        m = "balanced"
    g["mode"] = m
    memory.long["gamification"] = g
    memory._save()
    return m

# ------------- TTS (single speak implementation) -------------
# On servers (Render/containers), avoid initializing pyttsx3 (needs eSpeak)
engine = None
if pyttsx3 and not ON_SERVER:
    try:
        engine = pyttsx3.init()
    except Exception:
        engine = None

def emit_to_ui(key: str, payload: dict):
    try:
        socketio.emit(key, payload, namespace="/", broadcast=True)
    except Exception as e:
        print("Emit error:", e)

# Reconfirm TTS availability with safe init (local only)
if pyttsx3 and not ON_SERVER and engine is None:
    try:
        engine = pyttsx3.init()
    except Exception as e:
        engine = None
        print("TTS disabled:", e)

def speak(text):
    print("Draco says:", text)
    if engine:
        try:
            engine.say(text)
            engine.runAndWait()
        except Exception as e:
            print("TTS error:", e)

# ------------- Notes / Reminders -------------
class NotesManager:
    def __init__(self, path=NOTES_FILE):
        self.path = path
        self.notes = safe_read_json(self.path, [])

    def add(self, text):
        nid = int(time.time() * 1000)
        self.notes.append({"id": nid, "text": text, "created": datetime.datetime.now().isoformat()})
        safe_write_json(self.path, self.notes)
        return nid

    def delete(self, nid):
        self.notes = [n for n in self.notes if n["id"] != nid]
        safe_write_json(self.path, self.notes)

    def list(self):
        return list(self.notes)

class ReminderManager:
    def __init__(self, path=REMINDERS_FILE):
        self.path = path
        self.reminders = safe_read_json(self.path, [])
        self._start_loop()

    def add(self, text, when: datetime.datetime):
        rid = int(time.time() * 1000)
        self.reminders.append({"id": rid, "text": text, "at": when.isoformat()})
        safe_write_json(self.path, self.reminders)
        return rid

    def remove(self, rid):
        self.reminders = [r for r in self.reminders if r["id"] != rid]
        safe_write_json(self.path, self.reminders)

    def list(self):
        return list(self.reminders)

    def _start_loop(self):
        t = threading.Thread(target=self._loop, daemon=True)
        t.start()

    def _loop(self):
        while True:
            now = datetime.datetime.now()

            to_fire = []
            for r in list(self.reminders):
                try:
                    at = datetime.datetime.fromisoformat(r["at"])
                    if now >= at:
                        to_fire.append(r)
                except Exception:
                    continue
            for r in to_fire:
                msg = f"Reminder: {r['text']}"
                speak(msg)
                try:
                    # actionable notification in chat + beep alarm
                    emit_to_ui("draco_response", {"text": msg})
                    emit_to_ui("reminder_fired", {"text": r["text"]})
                    emit_to_ui("play_beep", {"kind": "reminder"})
                except Exception:
                    pass
                self.remove(r["id"])
            time.sleep(6)


notes_mgr = NotesManager()
reminder_mgr = ReminderManager()

# ------------- Pomodoro / Focus Timers -------------

class PomodoroState:
    def __init__(self):
        self.active = False
        self.mode = None  # "pomodoro", "break", "mini", "focus"
        self.total_seconds = 0
        self.remaining_seconds = 0
        self.paused = False
        self.started_at = None


class PomodoroManager:
    def __init__(self):
        self.state = PomodoroState()
        self._lock = threading.Lock()
        self._thread = threading.Thread(target=self._loop, daemon=True)
        self._thread.start()

    def start(self, seconds: int, mode: str):
        with self._lock:
            self.state.active = True
            self.state.mode = mode
            self.state.total_seconds = max(1, int(seconds))
            self.state.remaining_seconds = self.state.total_seconds
            self.state.paused = False
            self.state.started_at = time.time()
        payload = {
            "mode": mode,
            "total_seconds": self.state.total_seconds,
            "remaining_seconds": self.state.remaining_seconds,
            "progress": 0.0,
        }
        emit_to_ui("pomodoro_started", payload)
        emit_to_ui("pomodoro_tick", payload)
        emit_to_ui("play_beep", {"kind": "pomodoro_start"})

    def stop(self):
        with self._lock:
            self.state = PomodoroState()
        emit_to_ui("pomodoro_stopped", {})

    def pause(self):
        with self._lock:
            if self.state.active:
                self.state.paused = True
        emit_to_ui("pomodoro_paused", {})

    def resume(self):
        with self._lock:
            if self.state.active and self.state.paused:
                self.state.paused = False
        emit_to_ui("pomodoro_resumed", {})

    def reset(self):
        with self._lock:
            if self.state.total_seconds > 0:
                total = self.state.total_seconds
                mode = self.state.mode or "pomodoro"
            else:
                total = 25 * 60
                mode = "pomodoro"
            self.state.active = True
            self.state.mode = mode
            self.state.remaining_seconds = total
            self.state.paused = False
            self.state.started_at = time.time()
        payload = {
            "mode": mode,
            "total_seconds": total,
            "remaining_seconds": total,
            "progress": 0.0,
        }
        emit_to_ui("pomodoro_reset", payload)
        emit_to_ui("pomodoro_tick", payload)

    def get_status(self):
        with self._lock:
            st = self.state
            return {
                "active": st.active,
                "mode": st.mode,
                "total_seconds": st.total_seconds,
                "remaining_seconds": st.remaining_seconds,
                "paused": st.paused,
            }

    def _loop(self):
        while True:
            time.sleep(1)

            with self._lock:
                st = self.state
                if not st.active or st.paused or st.remaining_seconds <= 0:
                    continue
                st.remaining_seconds -= 1
                remaining = st.remaining_seconds
                total = st.total_seconds or 1
                mode = st.mode or "pomodoro"

            # emit tick update outside lock
            progress = max(0.0, min(1.0, 1.0 - (remaining / float(total))))
            socketio.emit(
                "pomodoro_tick",
                {
                    "mode": mode,
                    "remaining_seconds": remaining,
                    "total_seconds": total,
                    "progress": progress,
                },
                namespace="/",
                broadcast=True,
            )
            emit_to_ui("play_beep", {"kind": "tick"})
            if remaining <= 0:
                final_text = "Pomodoro done! ☕ All gone! Take a short break 😎"
                speak(final_text)
                # push a chat message plus dedicated events for UI/sound
                emit_to_ui("draco_response", {"text": final_text})
                emit_to_ui("pomodoro_done", {"mode": mode})
                emit_to_ui("play_beep", {"kind": "pomodoro_end"})
                with self._lock:
                    self.state = PomodoroState()


pomodoro_mgr = PomodoroManager()

# ------------- System utilities -------------
def system_status_summary():
    if not psutil:
        return "psutil not installed, cannot fetch system info."
    cpu = psutil.cpu_percent(interval=1)
    ram = psutil.virtual_memory()
    disk = psutil.disk_usage(os.path.expanduser("~"))
    bat = psutil.sensors_battery() if hasattr(psutil, "sensors_battery") else None
    lines = [
        f"CPU: {cpu}%",
        f"RAM: {ram.percent}%",
        f"Disk: {disk.percent}%"
    ]
    if bat:
        lines.append(f"Battery: {bat.percent}% {'charging' if bat.power_plugged else 'not charging'}")
    return ". ".join(lines)

def take_screenshot(save_path: Optional[str] = None):
    if not pyautogui:
        raise RuntimeError("pyautogui not installed.")
    if not save_path:
        save_path = f"screenshot_{int(time.time())}.png"
    img = pyautogui.screenshot()
    img.save(save_path)
    return save_path

def open_app_windows(name: str):
    """Try to open an app by name (Windows start) or command on other OS."""
    try:
        if platform.system() == "Windows":
            # try known apps mapping first
            mapping = {
                "vscode": r"C:\Users\%USERNAME%\AppData\Local\Programs\Microsoft VS Code\Code.exe",
                "spotify": r"C:\Users\%USERNAME%\AppData\Roaming\Spotify\Spotify.exe",
                "chrome": r"chrome",
            }
            if name.lower() in mapping:
                path = mapping[name.lower()]
                try:
                    os.startfile(os.path.expandvars(path))
                    return True
                except Exception:
                    # fallback to shell start
                    subprocess.run(["start", path], shell=True)
                    return True
            else:
                # best-effort start name
                subprocess.run(["start", name], shell=True)
                return True
        elif platform.system() == "Linux":
            subprocess.Popen([name])
            return True
        elif platform.system() == "Darwin":
            subprocess.Popen(["open", "-a", name])
            return True
    except Exception as e:
        print("open_app error:", e)
    return False

def run_system_command(cmd: str):
    """Run a shell command and return (stdout, stderr)."""
    try:
        res = subprocess.run(cmd, shell=True, capture_output=True, text=True, timeout=20)
        return res.stdout.strip(), res.stderr.strip()
    except Exception as e:
        return "", str(e)

# Brightness / wifi / bluetooth / volume - best-effort examples (Windows)
def set_brightness_win(percent: int):
    try:
        script = f'(Get-WmiObject -Namespace root/WMI -Class WmiMonitorBrightnessMethods).WmiSetBrightness(1,{int(percent)})'
        subprocess.run(["powershell", "-Command", script], check=False)
        return True
    except Exception as e:
        print("brightness error:", e)
        return False

def toggle_wifi_win(enable: bool):
    try:
        state = "enable" if enable else "disable"
        subprocess.run(["powershell", "-Command", f"Get-NetAdapter | Where-Object {{ $_.Status -ne '{state}' }} | ForEach-Object {{ Set-NetAdapter -Name $_.Name -{state} -Confirm:$false }}"], check=False)
        return True
    except Exception as e:
        print("wifi toggle error:", e)
        return False

# ------------- Web helpers -------------
def open_youtube():
    webbrowser.open("https://youtube.com")
    return "Opened YouTube."

def open_instagram():
    webbrowser.open("https://instagram.com")
    return "Opened Instagram."

def open_github():
    webbrowser.open("https://github.com")
    return "Opened GitHub."

def open_render():
    webbrowser.open("https://render.com")
    return "Opened Render."

def open_whatsapp_web():
    webbrowser.open("https://web.whatsapp.com")
    return "Opened WhatsApp Web."

def send_whatsapp_message(phone: str, message: str):
    """Open WhatsApp web chat to phone with message (uses wa.me). Works if user is logged-in."""
    # phone should be in international format without +, e.g., 919xxxxxxxxx
    msg = url_quote(message)
    url = f"https://wa.me/{phone}?text={msg}"
    if ON_SERVER:
        return {"text": f"Opening WhatsApp chat for {phone}…", "action": "open_url", "url": url}
    webbrowser.open(url)
    return f"Opening WhatsApp chat for {phone}. Please confirm send in browser."

def open_linkedin():
    webbrowser.open("https://linkedin.com")
    return "Opened LinkedIn."

# ------------- Music integration -------------
def play_music_from_library(song_name: Optional[str] = None):
    if not musicLibrary:
        return "musicLibrary not found. Add musicLibrary.py with play/pause/stop functions."
    # On server: return a web action so the browser opens the URL
    if ON_SERVER:
        try:
            url = None
            song_display_name = song_name or "music"
            # Use musicLibrary's built-in search function if available
            if hasattr(musicLibrary, "_find_song_url_by_name") and song_name:
                url = musicLibrary._find_song_url_by_name(song_name)
                if url:
                    return {"text": f"Playing: {song_name}", "action": "open_url", "url": url}
            # Fallback: search music dict directly
            if hasattr(musicLibrary, "music") and isinstance(musicLibrary.music, dict):
                if song_name:
                    key = song_name.strip().lower()
                    # Exact match
                return {"text": f"Playing: {song_display_name}", "action": "open_url", "url": url}
            return "No matching song found. Try: play faded, play lily, play alone"
        except Exception as e:
            return f"musicLibrary error: {e}"
    # Local: invoke the library which opens the browser
    try:
        if song_name:
            return musicLibrary.play(song_name)
        else:
            return musicLibrary.play()
    except Exception as e:
        return f"musicLibrary error: {e}"

# ------------- Local system actions (best-effort) -------------
def sleep_pc():
    if platform.system() == "Windows":
        try:
            subprocess.run(["rundll32.exe", "powrprof.dll,SetSuspendState", "0,1,0"], check=False)
            return True, "Sleeping now."
        except Exception as e:
            return False, f"Sleep failed: {e}"
    return False, "Sleep not supported on this OS."

def type_text(text: str):
    if not pyautogui:
        return False, "Typing requires pyautogui installed locally."
    try:
        pyautogui.typewrite(text, interval=0.02)
        return True, "Typed your text."
    except Exception as e:
        return False, f"Typing failed: {e}"

def set_brightness(percent: int):
    if platform.system() == "Windows":
        ok = set_brightness_win(percent)
        return ok, ("Brightness updated." if ok else "Couldn't change brightness.")
    return False, "Brightness control not supported on this OS."

def toggle_wifi(enable: bool):
    if platform.system() == "Windows":
        ok = toggle_wifi_win(enable)
        return ok, ("Wi-Fi enabled." if (ok and enable) else ("Wi-Fi disabled." if ok else "Couldn't toggle Wi-Fi."))
    return False, "Wi-Fi control not supported on this OS."

def set_volume(percent: int):
    return False, "Volume control not implemented on this system."

def toggle_bluetooth(enable: bool):
    return False, "Bluetooth control not implemented."

def whatsapp_send_direct(phone: str, message: str):
    if ON_SERVER:
        return False, "Direct WhatsApp send only works on your local machine."
    if not pywhatkit:
        return False, "Install pywhatkit locally to send WhatsApp instantly."
    try:
        if not phone.startswith("+") and len(phone) >= 10:
            phone = "+" + phone
        pywhatkit.sendwhatmsg_instantly(phone_no=phone, message=message, wait_time=8, tab_close=True)
        return True, "WhatsApp message sent."
    except Exception as e:
        return False, f"WhatsApp send failed: {e}"

WEATHER_API_KEY = "96e102262e9745129d110636251111 "

def get_weather(city):
    if not WEATHER_API_KEY:
        return "Weather feature needs an API key. Add it to WEATHER_API_KEY."
    
    try:
        # OpenWeatherMap API call
        url = f"http://api.openweathermap.org/data/2.5/weather?q={city}&appid={WEATHER_API_KEY}&units=metric"
        response = requests.get(url)
        data = response.json()

        if data.get("cod") != 200:
            return f"Could not find weather for '{city}'. Please check the city name."
        
        # Extract info
        temp = data["main"]["temp"]
        condition = data["weather"][0]["description"].title()
        humidity = data["main"]["humidity"]
        wind_speed = data["wind"]["speed"]

        # Format result
        result = (
            f"Weather in {city}:\n"
            f"Temperature: {temp}°C\n"
            f"Condition: {condition}\n"
            f"Humidity: {humidity}%\n"
            f"Wind Speed: {wind_speed} m/s"
        )
        return result

    except Exception as e:
        return f"Error fetching weather: {str(e)}"

NEWS_API_KEY = "231a93af4a8b4dc6bafbb736c20b20c3"

def get_news(topic="general"):
    if not NEWS_API_KEY:
        return "News feature needs an API key. Add it to NEWS_API_KEY."
    
    try:
        url = f"https://newsapi.org/v2/top-headlines?q={topic}&apiKey={NEWS_API_KEY}&pageSize=5"
        response = requests.get(url)
        data = response.json()
        
        if data.get("status") != "ok" or not data.get("articles"):
            return f"No news found for '{topic}'."
        
        # Format top 5 headlines
        news_list = ""
        for i, article in enumerate(data["articles"], start=1):
            title = article.get("title", "No title")
            source = article.get("source", {}).get("name", "Unknown")
            news_list += f"{i}. {title} ({source})\n"
        
        return f"Top news on '{topic}':\n{news_list}"

    except Exception as e:
        return f"Error fetching news: {str(e)}"

def solve_math(cmd):
    try:
        if "calculate" in cmd.lower():
            expression = cmd.lower().replace("calculate", "").strip()
            result = eval(expression)  # simple math
            speak(f"Result: {result}")
            return f"Result: {result}"

        elif "solve" in cmd.lower():
            equation = cmd.lower().replace("solve", "").strip()
            x = sp.symbols('x')
            solution = sp.solve(equation, x)
            speak(f"Solution: {solution}")
            return f"Solution: {solution}"

        else:
            speak("Math command not recognized.")
            return "Math command not recognized."

    except Exception as e:
        speak(f"Error: {str(e)}")
        return f"Error: {str(e)}"

EXCHANGE_API_KEY = "d97d653e87f3ea812b311d20"  # for currency

def convert_unit(cmd):
    try:
        cmd_lower = cmd.lower()

        # Currency Conversion
        if "convert" in cmd_lower and "to" in cmd_lower:
            words = cmd_lower.replace("convert", "").strip().split(" ")
            amount = float(words[0])
            from_unit = words[1].upper()
            to_unit = words[-1].upper()

            # Only handle USD/INR for free example, expand later
            if from_unit == "USD" and to_unit == "INR":
                rate = 82.5  # example rate
                result = amount * rate
                speak(f"{amount} {from_unit} = {result} {to_unit}")
                return f"{amount} {from_unit} = {result} {to_unit}"
            else:
                speak(f"Conversion from {from_unit} to {to_unit} not supported yet.")
                return f"Conversion from {from_unit} to {to_unit} not supported yet."

        # Simple Length Conversion Example
        elif "km to miles" in cmd_lower:
            km = float(cmd_lower.split("km")[0].strip())
            miles = km * 0.621371
            speak(f"{km} km = {miles:.2f} miles")
            return f"{km} km = {miles:.2f} miles"

        # Temperature Conversion Example
        elif "c to f" in cmd_lower:
            c = float(cmd_lower.split("c")[0].strip())
            f = (c * 9/5) + 32
            speak(f"{c}°C = {f:.2f}°F")
            return f"{c}°C = {f:.2f}°F"

        else:
            speak("Unit conversion not recognized.")
            return "Unit conversion not recognized."

    except Exception as e:
        speak(f"Error: {str(e)}")
        return f"Error: {str(e)}"

# ------------- Search & utilities -------------
def web_search_duckduckgo(query: str, limit: int = 3):
    """Return structured DuckDuckGo search results (best-effort)."""
    query = (query or "").strip()
    if not query:
        return []
    if DDGS is None:
        return []
    results = []
    try:
        with DDGS() as ddgs:
            for r in ddgs.text(query, max_results=max(1, limit)):
                title = (r.get("title") or r.get("href") or "Result").strip()
                snippet = (r.get("body") or "").strip()
                href = (r.get("href") or "").strip()
                results.append({
                    "title": title,
                    "snippet": snippet,
                    "href": href,
                })
                if len(results) >= limit:
                    break
    except Exception:
        return []
    return results


SMALL_TALK_RESPONSES = {
    "hi": "Hey there! ",
    "hello": "Hello! Ready when you are.",
    "how are you": "Doing great and ready to help!",
    "thanks": "Happy to help!",
    "thank you": "Anytime!",
    "good morning": "Good morning! Let's make it productive.",
    "good night": "Good night! Rest well.",
}


def _handle_small_talk(cmd_lower: str):
    for key, reply in SMALL_TALK_RESPONSES.items():
        if key in cmd_lower:
            speak(reply)
            return reply
    return None


def process_command(raw_cmd: str):
    """Map raw user phrases to actions."""
    if not raw_cmd:
        return "Please say something."

    cmd_clean = raw_cmd.strip()
    cmd = cmd_clean.lower()

    _maybe_award_keyword_xp(cmd)

    memory.add(f"You: {raw_cmd}")
    user_email = get_logged_in_email()
    if user_email:
        try:
            save_chat_line(user_email, "user", raw_cmd)
        except Exception:
            pass

    personality.update(cmd)

    # Small talk first
    small = _handle_small_talk(cmd)
    if small:
        return small

    # Intro handler
    if _handle_intro_request:
        intro_reply = _handle_intro_request(cmd)
        if intro_reply:
            speak(intro_reply)
            return intro_reply

    # Reminder block
    if "remind me to" in cmd or "set reminder to" in cmd:
        try:
            if " at " in cmd:
                before, atpart = cmd.split(" at ", 1)
            elif "at" in cmd:
                before, atpart = cmd.split("at", 1)
            else:
                return "Please include time with 'at'. Example: remind me to call mom at 19:30"
            text = (
                before.replace("remind me to", "")
                .replace("set reminder to", "")
                .strip()
            )
            if not text:
                return "What should I remind you about?"
            now = datetime.datetime.now()
            atpart = atpart.strip()
            if ":" in atpart and len(atpart.split()) == 1:
                hour, minute = map(int, atpart.split(":"))
                when = now.replace(hour=hour, minute=minute, second=0, microsecond=0)
                if when < now:
                    when += datetime.timedelta(days=1)
            else:
                when = datetime.datetime.fromisoformat(atpart)
            rid = reminder_mgr.add(text, when)
            r = f"Reminder set for {when.isoformat()} (ID {rid})"
            speak(r)
            emit_to_ui("reminder_created", {"id": rid, "text": text, "when": when.isoformat()})
            return r
        except Exception as e:
            return f"Couldn't set reminder: {e}"

    if "list reminders" in cmd:
        data = reminder_mgr.list()
        if not data:
            return "No reminders saved."
        lines = [f"{r['id']}: {r['text']} at {r['at']}" for r in data]
        return "Reminders:\n" + "\n".join(lines)

    if cmd.startswith("delete reminder"):
        try:
            parts = cmd.split()
            rid = int(parts[-1])
            reminder_mgr.remove(rid)
            return f"Deleted reminder {rid}."
        except Exception:
            return "Please say: delete reminder <id>."

    # Notes
    if cmd.startswith("take note") or cmd.startswith("note "):
        original = cmd_clean
        if cmd.startswith("take note"):
            note_text = original[len("take note"):].strip()
        else:
            note_text = original[len("note "):].strip()
        if not note_text:
            return "Please say the note text."
        nid = notes_mgr.add(note_text)
        add_xp("study", 4)
        r = f"Note saved. ID {nid}"
        speak(r)
        return r

    if "list notes" in cmd or "show notes" in cmd:
        notes = notes_mgr.list()
        if not notes:
            return "No notes."
        short = "; ".join([f"{i + 1}. {n['text']}" for i, n in enumerate(notes[:6])])
        speak("Reading notes.")
        add_xp("study", 2)
        return short

    search_query = _extract_search_query_text(cmd_clean, cmd)
    if search_query is not None:
        if not search_query:
            return "What should I search for?"
        _, summary = run_structured_search(search_query, limit=5, source="command")
        add_xp("study", 6)
        speak(f"Sharing top search hits for {search_query}.")
        return summary

    # Content generation / academics
    if cmd.startswith("generate doc on ") or cmd.startswith("create doc on "):
        topic = cmd.replace("generate doc on ", "", 1).replace("create doc on ", "", 1).strip()
        if not topic:
            return "Please provide a topic for the document."
        points = research_query_to_texts(topic, limit=10)
        try:
            path = _generate_docx(f"{topic.title()} - Notes", points)
            rel = os.path.relpath(path, os.getcwd()).replace("\\", "/")
            url = f"/download/{rel}"
            speak("Your document is ready, I've created a download link.")
            return {"text": f"Generated DOCX for {topic}. Download: {url}", "action": "open_url", "url": url}
        except Exception as e:
            return f"Could not generate DOCX: {e}"

    if cmd.startswith("generate pdf on ") or cmd.startswith("create pdf on "):
        topic = cmd.replace("generate pdf on ", "", 1).replace("create pdf on ", "", 1).strip()
        if not topic:
            return "Please provide a topic for the PDF."
        points, sources = research_query_to_texts_with_sources(topic, limit=12)
        try:
            path = _generate_pdf(f"{topic.title()} - Report", points, sources=sources)
            rel = os.path.relpath(path, os.getcwd()).replace("\\", "/")
            url = f"/download/{rel}"
            speak("Your PDF report is ready, I've created a download link.")
            return {"text": f"Generated PDF for {topic}. Download: {url}", "action": "open_url", "url": url}
        except Exception as e:
            return f"Could not generate PDF: {e}"

    if cmd.startswith("generate notes on ") or cmd.startswith("create notes on "):
        topic = cmd.replace("generate notes on ", "", 1).replace("create notes on ", "", 1).strip()
        if not topic:
            return "Please provide a topic for the notes."
        points, sources = research_query_to_texts_with_sources(topic, limit=10)
        try:
            lines = list(points)
            if sources:
                lines.append("")
                lines.extend(sources)
            speak("Your notes are ready.")
            return {"text": f"Generated notes for {topic}.", "action": "show_notes", "notes": lines}
        except Exception as e:
            return f"Could not generate notes: {e}"

    # Summarize recent chat (study helper)
    if "summarize this chat" in cmd or "chat summary" in cmd:
        sess = memory.get_session()[-40:]
        if not sess:
            return "There isn't enough recent chat to summarize yet."
        text_blob = "\n".join(x.get("text", "") for x in sess if x.get("text"))
        summary = _summarize_text(text_blob, max_len=800)
        add_xp("study", 10)
        return "Here is a short summary of our recent chat:\n" + summary

    # Generate practice questions from recent context
    if "generate questions" in cmd or "practice questions" in cmd:
        sess = memory.get_session()[-40:]
        if not sess:
            return "I need some recent content or explanation to turn into questions. Try asking a concept first."
        text_blob = "\n".join(x.get("text", "") for x in sess if x.get("text"))
        pts = _extract_key_points(text_blob, limit=8)
        if not pts:
            return "I couldn't find enough material to build questions from. Try explaining the topic again."
        lines = []
        for i, p in enumerate(pts, 1):
            base = p.strip().rstrip(".?!")
            q = base
            if not q.lower().startswith(("what", "why", "how", "when", "where", "who")):
                q = "What about: " + q
            lines.append(f"Q{i}. {q}?")
        add_xp("study", 12)
        return "Here are some practice questions based on our recent discussion:\n" + "\n".join(lines)

    # One-step coding practice questions by language
    if "python practice question" in cmd:
        add_xp("coding", 10)
        qs = [
            "Write a Python function that returns True if a string is a palindrome.",
            "Given a list of numbers, return the second largest element.",
            "Explain the difference between a list, tuple, and set in Python.",
        ]
        return "Here are some Python practice questions:\n- " + "\n- ".join(qs)

    if "c++ practice question" in cmd or "cpp practice question" in cmd:
        add_xp("coding", 10)
        qs = [
            "Write a C++ program to reverse an array in-place.",
            "What is a reference and how is it different from a pointer in C++?",
            "Implement a simple class for a BankAccount with deposit and withdraw methods.",
        ]
        return "Here are some C++ practice questions:\n- " + "\n- ".join(qs)

    if "java practice question" in cmd:
        add_xp("coding", 10)
        qs = [
            "Explain the difference between an interface and an abstract class in Java.",
            "Write a Java method to check if a number is prime.",
            "What is the purpose of the 'static' keyword in Java?",
        ]
        return "Here are some Java practice questions:\n- " + "\n- ".join(qs)

    if "javascript practice question" in cmd or "js practice question" in cmd:
        add_xp("coding", 10)
        qs = [
            "What is the difference between 'let', 'const', and 'var' in JavaScript?",
            "Write a function that debounces another function.",
            "Explain how promises work and what 'async/await' does.",
        ]
        return "Here are some JavaScript practice questions:\n- " + "\n- ".join(qs)

    if "html practice question" in cmd or "css practice question" in cmd:
        add_xp("coding", 10)
        qs = [
            "Create a simple HTML page with a header, footer, and a main section.",
            "Write CSS to center a div both vertically and horizontally.",
            "Explain the difference between inline, inline-block, and block elements.",
        ]
        return "Here are some HTML/CSS practice questions:\n- " + "\n- ".join(qs)

    # One-step quick coding tips by language
    if "quick python tip" in cmd or "python quick tip" in cmd:
        add_xp("coding", 5)
        tip = "Use list comprehensions and 'enumerate' to write clean loops, and always prefer 'with open(...)' for file handling."
        return f"Quick Python tip: {tip}"

    if "quick c++ tip" in cmd or "quick cpp tip" in cmd:
        add_xp("coding", 5)
        tip = "Prefer std::vector over raw arrays, and initialize variables using brace initialization to avoid surprises."
        return f"Quick C++ tip: {tip}"

    if "quick java tip" in cmd:
        add_xp("coding", 5)
        tip = "Keep your classes small and focused, and always program to interfaces rather than concrete implementations."
        return f"Quick Java tip: {tip}"

    if "quick javascript tip" in cmd or "quick js tip" in cmd:
        add_xp("coding", 5)
        tip = "Avoid global variables, use 'const' and 'let', and keep async code readable with async/await."
        return f"Quick JavaScript tip: {tip}"

    if "quick html tip" in cmd or "quick css tip" in cmd:
        add_xp("coding", 5)
        tip = "Use semantic HTML tags and keep your CSS modular with utility classes or BEM-style naming."
        return f"Quick HTML/CSS tip: {tip}"

    # Coding helpers (non-LLM heuristics)
    if "explain this code" in cmd or cmd.startswith("explain code"):
        add_xp("coding", 10)
        return (
            "Here's how I approach explaining code without running it:\n"
            "1. Break it into functions and blocks.\n"
            "2. Identify inputs, outputs, and side effects.\n"
            "3. Track the main data structures through the code path.\n"
            "Paste the code here and I will walk you through it step by step."
        )

    if "find bug" in cmd or "debug this code" in cmd or "debug my code" in cmd:
        add_xp("coding", 12)
        return (
            "To debug this code together, we'll do this:\n"
            "- Add print/log statements around the part that fails.\n"
            "- Check variable values against what you expect.\n"
            "- Narrow down the smallest snippet that still shows the bug.\n"
            "Send me the error message and the code, and I'll help you reason through it."
        )

    if "refactor this" in cmd or "refactor code" in cmd:
        add_xp("coding", 10)
        return (
            "Refactoring plan:\n"
            "1. Extract repeated logic into small helper functions.\n"
            "2. Rename variables and functions to be more descriptive.\n"
            "3. Separate input/processing/output into clear layers.\n"
            "Paste the code and I will suggest a cleaner structure."
        )

    # Run shell command
    if cmd.startswith("run "):
        to_run = cmd.replace("run ", "", 1)
        out, err = run_system_command(to_run)
        if out:
            speak("Command executed, returning output.")
            return out[:1500]
        return f"Command error: {err}"

    # Final fallback
    speak("I didn't get that. Try asking me to open apps, play music, take notes, set reminders or search the web.")
    return "Unknown command. Try: open youtube, play music, take note, set reminder, search for ..."
    # Always show main app; login removed
    return send_from_directory(".", "draco.html")

@app.route("/guest")
def guest_mode():
    # Guest mode (login removed) – just serve app
    return send_from_directory(".", "draco.html")

@app.route("/api/gamestate", methods=["GET"])
def api_gamestate():
    """Return current gamification state for the XP bar and mode chips."""
    try:
        g = get_gamestate()
        return {"ok": True, "gamification": g}
    except Exception as e:
        return {"ok": False, "error": str(e)}, 500

@app.route("/api/gamestate/mode", methods=["POST"])
def api_gamestate_mode():
    """Update the active gamification mode (study / coding / balanced)."""
    data = request.json or {}
    mode = data.get("mode", "")
    try:
        m = set_mode(mode)
        g = get_gamestate()
        return {"ok": True, "mode": m, "gamification": g}
    except Exception as e:
        return {"ok": False, "error": str(e)}, 500

@app.route("/api/profile", methods=["GET", "POST"])
def api_profile():
    email = get_logged_in_email()
    if request.method == "GET":
        return {"ok": True, "profile": get_user_profile(email)}
    data = request.json or {}
    if not isinstance(data, dict):
        return {"ok": False, "error": "invalid"}, 400
    prof = get_user_profile(email)
    prof.update({k: v for k, v in data.items() if isinstance(k, str)})
    set_user_profile(email, prof)
    return {"ok": True, "profile": prof}

@app.route("/api/profile/clear", methods=["POST"])
def api_profile_clear():
    email = get_logged_in_email()
    set_user_profile(email, {})
    return {"ok": True}

@app.route("/api/guest_profile", methods=["GET", "POST"])
def api_guest_profile():
    # Use MemoryManager.long as guest profile store
    if request.method == "GET":
        return {"ok": True, "profile": memory.long}
    data = request.json or {}
    if not isinstance(data, dict):
        return {"ok": False, "error": "invalid"}, 400
    try:
        # Only accept simple keys to avoid arbitrary writes
        allowed = {"name", "hobbies", "favorite_subject"}
        for k, v in data.items():
            if k in allowed:
                memory.set_pref(k, v)
        return {"ok": True, "profile": memory.long}
    except Exception as e:
        return {"ok": False, "error": str(e)}, 500

@app.route("/api/guest_profile/clear", methods=["POST"])
def api_guest_profile_clear():
    try:
        for k in ["name", "hobbies", "favorite_subject"]:
            if k in memory.long:
                del memory.long[k]
        memory._save()
        return {"ok": True}
    except Exception as e:
        return {"ok": False, "error": str(e)}, 500

@app.route("/api/chat_history", methods=["GET"]) 
def api_chat_history():
    email = get_logged_in_email()
    chat_id = request.args.get("chat_id")
    return {"ok": True, "items": get_chat_history(email, chat_id)}

@app.route("/api/chats", methods=["GET"]) 
def api_chats_list():
    email = get_logged_in_email()
    chats = _load_chats(email)
    # minimal list
    out = [
        {"id": c.get("id"), "name": c.get("name") or "New Chat", "updated_at": c.get("updated_at", 0)}
        for c in chats
    ]
    out.sort(key=lambda x: x.get("updated_at", 0), reverse=True)
    return {"ok": True, "chats": out, "current": _current_chat_id()}

@app.route("/api/chats/new", methods=["POST"]) 
def api_chats_new():
    email = get_logged_in_email()
    c = _create_new_chat(email)
    return {"ok": True, "chat": {"id": c.get("id"), "name": c.get("name")}}

@app.route("/api/chats/select", methods=["POST"]) 
def api_chats_select():
    email = get_logged_in_email()
    data = request.json or {}
    cid = str(data.get("chat_id", ""))
    if not cid:
        return {"ok": False, "error": "invalid"}, 400
    chats = _load_chats(email)
    if any(c.get("id") == cid for c in chats):
        _set_current_chat_id(cid)
        return {"ok": True}
    return {"ok": False, "error": "not_found"}, 404

@app.route("/api/chats/clear", methods=["POST"]) 
def api_chats_clear():
    email = get_logged_in_email()
    ok = _clear_current_chat(email)
    return {"ok": ok}

# ------------- Research & DOCX export -------------
GENERATED_DIR = os.path.join(os.getcwd(), "generated")
ensure_dir(GENERATED_DIR)
UPLOADS_DIR = os.path.join(os.getcwd(), "uploads")
ensure_dir(UPLOADS_DIR)

def research_query_to_texts(query: str, limit: int = 6):
    text = web_search_duckduckgo(query, limit=limit)
    if isinstance(text, str):
        parts = [p.strip() for p in text.split("|") if p.strip()]
        return parts[:limit] if parts else [text]
    return [str(text)]

def research_query_to_texts_with_sources(query: str, limit: int = 6):
    """Return (texts, sources_urls) using DDGS when available.
    Falls back to research_query_to_texts without sources.
    """
    if DDGS is None:
        return research_query_to_texts(query, limit=limit), []
    texts = []
    urls = []
    try:
        with DDGS() as ddgs:
            for r in ddgs.text(query, max_results=limit):
                body = (r.get("body") or "").strip()
                href = (r.get("href") or "").strip()
                if body:
                    # de-duplicate by normalized snippet
                    norm = " ".join(body.lower().split())
                    if all(" ".join(t.lower().split()) != norm for t in texts):
                        texts.append(body)
                if href:
                    urls.append(href)
        # ensure bounded sizes
        return texts[:limit] if texts else ["No results found."], urls[:limit]
    except Exception:
        return research_query_to_texts(query, limit=limit), []

def save_docx_from_texts(title: str, bullets):
    if not Document:
        return None, "python-docx not installed."
    try:
        doc = Document()
        doc.add_heading(title.strip() or "Research", level=1)
        doc.add_paragraph("")
        for b in bullets:
            doc.add_paragraph(b.strip(), style=None)
        safe_name = "".join(ch for ch in title if ch.isalnum() or ch in (" ", "_", "-")).strip() or "research"
        filename = f"{safe_name[:40].replace(' ', '_')}_{int(time.time())}.docx"
        path = os.path.join(GENERATED_DIR, filename)
        doc.save(path)
        return path, None
    except Exception as e:
        return None, str(e)

def _generate_docx(title: str, bullets):
    """Compatibility wrapper around save_docx_from_texts that returns only the path.

    Many parts of the codebase call _generate_docx(title, bullets) and expect a path.
    """
    path, err = save_docx_from_texts(title, bullets)
    if not path:
        raise RuntimeError(err or "failed to write DOCX")
    return path

@app.route("/api/research", methods=["POST"])
def api_research():
    """
    Body: { "query": "topic", "make_doc": true/false }
    Returns summaries and optional doc download path.
    """
    data = request.json or {}
    query = (data.get("query") or "").strip()
    make_doc = bool(data.get("make_doc"))
    if not query:
        return {"ok": False, "error": "missing_query"}, 400

    items = research_query_to_texts(query, limit=6)
    out = {"ok": True, "items": items}
    if make_doc:
        path, err = save_docx_from_texts(query, items)
        if path:
            # Provide relative path for download (served by Flask send_from_directory if needed)
            rel = os.path.relpath(path, os.getcwd()).replace("\\", "/")
            out["doc"] = f"/download/{rel}"
        else:
            out["doc_error"] = err or "failed_to_write_doc"
    return out

@app.route("/download/<path:filepath>", methods=["GET"])
def download_generated(filepath):
    # Only allow files inside GENERATED_DIR
    abs_path = os.path.abspath(os.path.join(os.getcwd(), filepath))
    if not abs_path.startswith(os.path.abspath(GENERATED_DIR)):
        return {"ok": False, "error": "forbidden"}, 403
    try:
        directory = os.path.dirname(abs_path)
        filename = os.path.basename(abs_path)
        return send_from_directory(directory, filename, as_attachment=True)
    except Exception as e:
        return {"ok": False, "error": str(e)}, 404

def _extract_text_from_docx(path: str) -> str:
    if not Document:
        return ""
    try:
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception:
        return ""

def _extract_text_from_pptx(path: str) -> str:
    if not Presentation:
        return ""
    try:
        prs = Presentation(path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    t = (shape.text or "").strip()
                    if t:
                        texts.append(t)
        return "\n".join(texts)
    except Exception:
        return ""

def _extract_text_from_pdf(path: str) -> str:
    if not PdfReader:
        return ""
    try:
        reader = PdfReader(path)
        texts = []
        for page in reader.pages:
            t = page.extract_text() or ""
            if t.strip():
                texts.append(t.strip())
        return "\n".join(texts)
    except Exception:
        return ""

def _summarize_text(text: str, max_len: int = 1200) -> str:
    # naive summarizer: keep first N chars and compress whitespace
    if not text:
        return "No content to summarize."
    s = " ".join(text.split())
    return s[:max_len] + ("…" if len(s) > max_len else "")

def _split_sentences(text: str):
    try:
        parts = re.split(r"(?<=[.!?])\s+", text.strip())
        return [p.strip() for p in parts if p and len(p.strip()) > 3]
    except Exception:
        return [text]

def _extract_key_points(text: str, limit: int = 8):
    sents = _split_sentences(text)
    # heuristic: prefer sentences containing key verbs and shorter than 200 chars
    keys = [" is ", " are ", " was ", " were ", " include ", " consists", " uses ", " means ", " defines ", " key ", " important "]
    scored = []
    for s in sents:
        score = 0
        ln = len(s)
        score += max(0, 200 - ln) / 50.0
        low = " " + s.lower() + " "
        score += sum(1 for k in keys if k in low)
        scored.append((score, s))
    scored.sort(key=lambda x: x[0], reverse=True)
    pts = [s for _, s in scored[:max(1, limit)]]
    return pts

def _make_flashcards(text: str, limit: int = 8):
    pts = _extract_key_points(text, limit=limit)
    cards = []
    for p in pts:
        pl = p.strip().rstrip(".?!")
        # Build simple Q from clause
        q = "What is: " + (pl[:80] + ("…" if len(pl) > 80 else ""))
        a = pl
        cards.append({"q": q, "a": a})
    return cards

def _outline_text(text: str, max_sections: int = 6):
    paras = [p.strip() for p in text.split("\n\n") if p.strip()]
    outline = []
    for i, p in enumerate(paras[:max_sections], 1):
        title = _split_sentences(p)[0] if _split_sentences(p) else p[:60]
        bullets = _extract_key_points(p, limit=5)
        outline.append({"title": f"Section {i}: " + title[:60], "bullets": bullets})
    return outline

def _clean_text(text: str):
    # normalize whitespace and remove duplicate consecutive lines
    t = "\n".join([ln.strip() for ln in text.replace("\r\n", "\n").replace("\r", "\n").split("\n")])
    lines = [ln for ln in t.split("\n") if ln]
    out = []
    last = None
    for ln in lines:
        if ln != last:
            out.append(ln)
            last = ln
    return "\n".join(out)

def _rewrite_tone(text: str, tone: str):
    # Simple heuristic rewrites; non-LLM
    t = text
    tone = (tone or "").lower()
    if tone == "simple":
        # shorter sentences
        sents = _split_sentences(t)
        sents = [s.split(',')[0].strip() for s in sents]
        return " ".join(sents)
    if tone == "formal":
        repl = {
            " can't ": " cannot ",
            " won't ": " will not ",
            " it's ": " it is ",
            " isn't ": " is not ",
            " don't ": " do not ",
            " doesn't ": " does not ",
            " i'm ": " I am ",
        }
        low = " " + t
        for k,v in repl.items():
            low = low.replace(k, v)
        return low.strip()
    if tone == "academic":
        # add connectors
        sents = _split_sentences(t)
        enh = []
        for i, s in enumerate(sents):
            prefix = "" if i == 0 else ("Furthermore, " if i % 2 else "Additionally, ")
            enh.append(prefix + s)
        return " ".join(enh)
    return t

def _extract_glossary(text: str, limit: int = 15):
    # heuristic: terms before ':' or capitalized words sequences
    terms = set()
    for ln in text.split("\n"):
        if ':' in ln:
            k = ln.split(':',1)[0].strip()
            if 2 <= len(k) <= 60:
                terms.add(k)
    words = re.findall(r"\b([A-Z][a-zA-Z]{2,}(?:\s+[A-Z][a-zA-Z]{2,})*)\b", text)
    for w in words:
        if 2 <= len(w) <= 40:
            terms.add(w)
    terms = list(terms)
    terms.sort(key=lambda x: x.lower())
    items = [f"{t}: " for t in terms[:limit]]
    return items

def _generate_pptx(title: str, bullets, max_sentences_per_slide: int = 4) -> str:
    if not Presentation:
        raise RuntimeError("python-pptx not installed.")
    prs = Presentation()

    # palettes and fonts for a modern look
    palettes = [
        {"bg": RGBColor(245, 248, 255), "accent": RGBColor(62, 99, 221)},  # soft blue
        {"bg": RGBColor(250, 247, 255), "accent": RGBColor(141, 82, 255)},  # soft violet
        {"bg": RGBColor(246, 252, 250), "accent": RGBColor(10, 163, 127)},  # teal
        {"bg": RGBColor(254, 248, 246), "accent": RGBColor(242, 95, 58)},   # coral
        {"bg": RGBColor(247, 249, 252), "accent": RGBColor(45, 55, 72)},    # slate
    ]
    fonts = ["Montserrat", "Poppins", "Roboto", "Segoe UI", "Arial"]
    palette = random.choice(palettes)
    font_name = random.choice(fonts)

    # helper to create a white rounded card
    def add_card(slide, margin=Inches(0.6)):
        left = margin
        top = margin
        width = prs.slide_width - margin * 2
        height = prs.slide_height - margin * 2
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.fill.background()
        return card

    # infer a simple icon from the title/topic and add title
    def get_topic_icon(title_text: str) -> str:
        t = (title_text or "").lower()
        mapping = [
            ("ai", "🤖"), ("machine learning", "🤖"), ("data", "📊"), ("analytics", "📈"),
            ("climate", "🌍"), ("environment", "🌿"), ("biology", "🧬"), ("health", "🩺"),
            ("medicine", "🧪"), ("finance", "💹"), ("marketing", "📣"), ("cloud", "☁️"),
            ("security", "🔒"), ("python", "🐍"), ("java", "☕"), ("web", "🌐"),
            ("quantum", "⚛️"), ("robot", "🤖"), ("network", "🌐"), ("database", "🗄️"),
        ]
        for k, e in mapping:
            if k in t:
                return e
        return "✨"

    def add_title(slide, text, y=Inches(1.0)):
        tx = slide.shapes.add_textbox(Inches(1.1), y, prs.slide_width - Inches(2.2), Inches(1.2))
        tf = tx.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = f"{get_topic_icon(title)}  {text}"
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.name = font_name
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.color.rgb = palette["accent"]
        return tx

    def add_icon_bullets(slide, count, top=Inches(2.1), left=Inches(0.95)):
        line_h = Inches(0.42)
        size = Inches(0.12)
        for idx in range(count):
            y = top + Inches(0.2) + line_h * idx
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, y, size, size)
            dot.fill.solid()
            dot.fill.fore_color.rgb = palette["accent"]
            dot.line.fill.background()
        return True

    def add_bullets(slide, items, top=Inches(2.1), left=Inches(1.1), width=None):
        if width is None:
            width = prs.slide_width - Inches(2.2)
        tx = slide.shapes.add_textbox(left, top, width, prs.slide_height - top - Inches(1.0))
        tf = tx.text_frame
        tf.clear()
        tf.word_wrap = True
        for idx, line in enumerate(items):
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.level = 0
            p.font.name = font_name
            p.font.size = Pt(22)
        add_icon_bullets(slide, len(items), top, left - Inches(0.15))
        return tx

    # Cover slide
    cover_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    cover = prs.slides.add_slide(cover_layout)
    # background
    try:
        cover.background.fill.solid()
        cover.background.fill.fore_color.rgb = palette["bg"]
    except Exception:
        pass
    add_card(cover, margin=Inches(0.7))
    add_title(cover, title, y=Inches(1.5))
    # subtitle
    sub = cover.shapes.add_textbox(Inches(1.1), Inches(2.3), prs.slide_width - Inches(2.2), Inches(0.8))
    stf = sub.text_frame
    stf.clear()
    sp = stf.paragraphs[0]
    sp.text = "Generated by Draco"
    sp.alignment = PP_ALIGN.LEFT
    srun = sp.runs[0]
    srun.font.name = font_name
    srun.font.size = Pt(20)

    # Content slides: split into sentences and batch into slides
    def split_sentences(text: str):
        raw = [s.strip() for s in re.split(r"(?<=[.!?۔؟！。]|।)\s+", text) if s.strip()]
        abbr = {"e.g.", "i.e.", "etc.", "mr.", "mrs.", "ms.", "dr.", "prof.", "sr.", "jr.", "vs.", "no.", "fig.", "al.", "u.s.", "u.k.", "dept.",
                "inc.", "ltd.", "co.", "est.", "approx.", "misc.", "ref.", "ed.", "pp.", "vol.", "jan.", "feb.", "mar.", "apr.", "jun.", "jul.", "aug.", "sep.", "oct.", "nov.", "dec."}
        merged = []
        for seg in raw:
            if merged and (merged[-1].lower().endswith(tuple(abbr)) or len(merged[-1]) <= 2 or re.search(r"\b[A-Z]\.\s*$", merged[-1])):
                merged[-1] = (merged[-1] + " " + seg).strip()
            else:
                merged.append(seg)
        return merged

    all_sentences = []
    for b in bullets:
        # support list of paragraphs/snippets
        all_sentences.extend(split_sentences(str(b)))

    if max_sentences_per_slide <= 0:
        max_sentences_per_slide = 4

    # batch sentences into slides
    for i in range(0, len(all_sentences), max_sentences_per_slide):
        chunk = all_sentences[i:i + max_sentences_per_slide]
        slide_palette = random.choice(palettes)
        layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[1]
        s = prs.slides.add_slide(layout)
        try:
            s.background.fill.solid()
            s.background.fill.fore_color.rgb = slide_palette["bg"]
        except Exception:
            pass
        variant = random.choice(["card", "accent_bar", "two_col", "image_text"])  # more variety
        if variant == "accent_bar":
            bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.35), prs.slide_height)
            bar.fill.solid()
            bar.fill.fore_color.rgb = slide_palette["accent"]
            bar.line.fill.background()
        card = add_card(s, margin=Inches(0.6))
        add_title(s, "Key Points", y=Inches(1.1))
        if variant == "two_col" and len(chunk) >= 2:
            mid = (len(chunk) + 1) // 2
            left_items = chunk[:mid]
            right_items = chunk[mid:]
            col_width = (prs.slide_width - Inches(2.2) - Inches(0.6)) / 2
            add_bullets(s, left_items, top=Inches(2.2), left=Inches(1.1), width=col_width)
            add_bullets(s, right_items, top=Inches(2.2), left=Inches(1.1) + col_width + Inches(0.6), width=col_width)
        elif variant == "image_text":
            # left image placeholder, right text
            img_left = Inches(1.1)
            img_top = Inches(2.2)
            img_w = Inches(3.6)
            img_h = Inches(3.0)
            ph = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, img_left, img_top, img_w, img_h)
            ph.fill.solid()
            ph.fill.fore_color.rgb = RGBColor(235, 240, 255)
            ph.line.color.rgb = slide_palette["accent"]
            txt = s.shapes.add_textbox(img_left, img_top + img_h + Inches(0.05), img_w, Inches(0.4))
            ttf = txt.text_frame
            ttf.clear()
            tp = ttf.paragraphs[0]
            tp.text = "Illustration"
            tp.font.name = font_name
            tp.font.size = Pt(12)
            # bullets on right
            add_bullets(s, chunk, top=Inches(2.2), left=img_left + img_w + Inches(0.6))
        else:
            add_bullets(s, chunk, top=Inches(2.2))

    safe = "".join(ch for ch in title if ch.isalnum() or ch in (" ","_","-")).strip() or "slides"
    name = f"{safe[:40].replace(' ','_')}_{int(time.time())}.pptx"
    path = os.path.join(GENERATED_DIR, name)
    prs.save(path)
    return path

def _generate_pdf(title: str, paragraphs, sources=None) -> str:
    if not FPDF:
        raise RuntimeError("fpdf not installed.")
    class PDFReport(FPDF):
        def __init__(self, t: str):
            super().__init__()
            self.t = t
            # theme colors (RGB)
            self.theme = random.choice([
                {"accent": (10, 163, 127), "muted": (100, 100, 120)},   # teal
                {"accent": (62, 99, 221), "muted": (110, 110, 140)},   # blue
                {"accent": (242, 95, 58), "muted": (130, 110, 110)},   # coral
                {"accent": (45, 55, 72), "muted": (110, 110, 120)},    # slate
            ])
            self.section_links = []  # list of (text, link_id)
        def header(self):
            try:
                self.set_font("Helvetica", "", 10)
                self.set_text_color(100)
                self.cell(0, 8, self.t, ln=1, align="C")
                self.set_draw_color(200)
                self.set_line_width(0.2)
                self.line(self.l_margin, self.get_y(), self.w - self.r_margin, self.get_y())
                self.ln(2)
            except Exception:
                pass
        def footer(self):
            try:
                self.set_y(-12)
                self.set_font("Helvetica", "", 9)
                self.set_text_color(120)
                self.cell(0, 10, f"Page {self.page_no()}", align="C")
            except Exception:
                pass
    pdf = PDFReport(title)
    pdf.set_margins(18, 16, 18)
    pdf.set_auto_page_break(auto=True, margin=18)
    try:
        pdf.set_title(title)
        pdf.set_author("Draco AI")
        pdf.set_subject("Generated Report")
    except Exception:
        pass
    pdf.add_page()
    pdf.set_text_color(0)
    pdf.set_font("Helvetica", "B", 18)
    pdf.cell(0, 10, txt=title, ln=True, align="C")
    pdf.ln(2)
    try:
        now_str = datetime.datetime.now().strftime("%b %d, %Y")
        pdf.set_font("Helvetica", "", 11)
        pdf.set_text_color(80)
        pdf.cell(0, 8, txt=now_str, ln=True, align="C")
    except Exception:
        pass
    pdf.ln(6)
    # Table of Contents placeholder page (will add after cover)
    # Create TOC after building section links; we'll add content links now and come back
    # Prepare section link ids
    section_link_ids = []
    for idx, p in enumerate(paragraphs, 1):
        try:
            link_id = pdf.add_link()
        except Exception:
            link_id = None
        section_link_ids.append((f"Section {idx}", link_id))

    # Add content pages with section headers and set link targets
    pdf.set_text_color(0)
    pdf.set_font("Helvetica", "", 12)
    for i, p in enumerate(paragraphs, 1):
        # Section header
        pdf.set_font("Helvetica", "B", 14)
        r,g,b = pdf.theme.get("accent", (0,0,0))
        pdf.set_text_color(r, g, b)
        header = f"Section {i}"
        y_before = pdf.get_y()
        pdf.cell(0, 8, txt=header, ln=True)
        # set link target for this section at current position
        try:
            if section_link_ids[i-1][1] is not None:
                pdf.set_link(section_link_ids[i-1][1], y=y_before, page=pdf.page_no())
        except Exception:
            pass
        pdf.set_text_color(0)
        pdf.set_font("Helvetica", "", 12)
        pdf.multi_cell(0, 7, txt=str(p))
        pdf.ln(2)

    # Insert Table of Contents page after the cover page
    try:
        pdf.set_auto_page_break(auto=True, margin=18)
        # Remember current state
        saved_page = pdf.page
        # Insert a page at position 2 (after cover)
        pdf.page = 1
        pdf._newpage("P")
        pdf.page = 2
        pdf.set_text_color(0)
        pdf.set_font("Helvetica", "B", 16)
        pdf.cell(0, 10, txt="Table of Contents", ln=True)
        pdf.ln(4)
        pdf.set_font("Helvetica", "", 12)
        for (label, link_id) in section_link_ids:
            if link_id is not None:
                pdf.set_text_color(0, 0, 180)
                pdf.write(8, label, link=link_id)
            else:
                pdf.set_text_color(0)
                pdf.cell(0, 8, txt=label, ln=True)
            pdf.ln(2)
        pdf.set_text_color(0)
        # restore to last page to continue sources and save
        pdf.page = saved_page + 1
    except Exception:
        # If insertion unsupported, just add TOC at the end
        try:
            pdf.add_page()
            pdf.set_text_color(0)
            pdf.set_font("Helvetica", "B", 16)
            pdf.cell(0, 10, txt="Table of Contents", ln=True)
            pdf.ln(4)
            pdf.set_font("Helvetica", "", 12)
            for (label, link_id) in section_link_ids:
                if link_id is not None:
                    pdf.set_text_color(0, 0, 180)
                    pdf.write(8, label, link=link_id)
                else:
                    pdf.set_text_color(0)
                    pdf.cell(0, 8, txt=label, ln=True)
                pdf.ln(2)
            pdf.set_text_color(0)
        except Exception:
            pass
    # Sources section (clickable links)
    if sources:
        pdf.ln(4)
        pdf.set_font("Helvetica", "B", 14)
        pdf.set_text_color(0)
        pdf.cell(0, 10, txt="Sources", ln=True)
        pdf.set_font("Helvetica", "", 11)
        for u in sources:
            try:
                display = u
                pdf.set_text_color(0, 0, 180)
                pdf.write(6, display, link=u)
                pdf.ln(6)
            except Exception:
                try:
                    pdf.set_text_color(0)
                    pdf.multi_cell(0, 6, txt=str(u))
                except Exception:
                    pass
        pdf.set_text_color(0)
    safe = "".join(ch for ch in title if ch.isalnum() or ch in (" ","_","-")).strip() or "report"
    name = f"{safe[:40].replace(' ','_')}_{int(time.time())}.pdf"
    path = os.path.join(GENERATED_DIR, name)
    pdf.output(path)
    return path

    

@app.route("/api/upload_process", methods=["POST"])
def api_upload_process():
    """
    Accept a user file and optional instruction.
    Returns processed summary or a modified file for download.
    Form fields:
      - file: uploaded file
      - instruction: optional text instructions (e.g., "summarize", "shorten", etc.)
    """
    if "file" not in request.files:
        return {"ok": False, "error": "no_file"}, 400
    f = request.files["file"]
    if f.filename == "":
        return {"ok": False, "error": "empty_filename"}, 400
    instruction = (request.form.get("instruction") or "").strip().lower()
    filename = secure_filename(f.filename)
    path = os.path.join(UPLOADS_DIR, filename)
    f.save(path)

    ext = os.path.splitext(filename)[1].lower()
    text = ""
    if ext == ".docx":
        text = _extract_text_from_docx(path)
    elif ext == ".pptx":
        text = _extract_text_from_pptx(path)
    elif ext == ".pdf":
        text = _extract_text_from_pdf(path)
    else:
        return {"ok": False, "error": "unsupported_type"}, 400

    if not text:
        return {"ok": False, "error": "no_text_found"}, 400

    # simple processing
    title = os.path.splitext(filename)[0]
    if "summarize" in instruction or "summary" in instruction:
        # Support presets like summarize:short|medium|detailed
        max_len = 1200
        try:
            if ":" in instruction:
                _, kind = instruction.split(":", 1)
                kind = (kind or "").strip().lower()
                if kind == "short":
                    max_len = 700
                elif kind == "medium":
                    max_len = 1200
                elif kind == "detailed":
                    max_len = 2200
        except Exception:
            pass
        summary = _summarize_text(text, max_len=max_len)
        try:
            out = _generate_docx(f"Summary of {title}", [summary])
            rel = os.path.relpath(out, os.getcwd()).replace("\\", "/")
            return {"ok": True, "summary": summary, "doc": f"/download/{rel}"}
        except Exception:
            return {"ok": True, "summary": summary}

    if "shorten" in instruction:
        short = _summarize_text(text, max_len=800)
        try:
            out = _generate_docx(f"Shortened - {title}", [short])
            rel = os.path.relpath(out, os.getcwd()).replace("\\", "/")
            return {"ok": True, "text": short, "doc": f"/download/{rel}"}
        except Exception:
            return {"ok": True, "text": short}

    if "lengthen" in instruction:
        augmented = text
        try:
            extra = web_search_duckduckgo(title, limit=3)
            if isinstance(extra, str) and extra and not extra.startswith("ddgs package not installed"):
                augmented = (text + "\n\n" + extra)[:8000]
        except Exception:
            pass
        try:
            out = _generate_docx(f"Extended - {title}", [augmented])
            rel = os.path.relpath(out, os.getcwd()).replace("\\", "/")
            return {"ok": True, "text": augmented[:3000], "doc": f"/download/{rel}"}
        except Exception:
            return {"ok": True, "text": augmented[:3000]}

    if "search" in instruction:
        try:
            results = web_search_duckduckgo(title, limit=5)
        except Exception as e:
            results = f"Search error: {e}"
        try:
            out = _generate_docx(f"Search Results - {title}", [str(results)])
            rel = os.path.relpath(out, os.getcwd()).replace("\\", "/")
            return {"ok": True, "text": str(results)[:3000], "doc": f"/download/{rel}"}
        except Exception:
            return {"ok": True, "text": str(results)[:3000]}

    if "keypoints" in instruction or "key points" in instruction:
        pts = _extract_key_points(text, limit=10)
        preview = "\n".join(f"- {p}" for p in pts)
        try:
            out = _generate_docx(f"Key Points - {title}", pts)
            rel = os.path.relpath(out, os.getcwd()).replace("\\", "/")
            return {"ok": True, "text": preview[:3000], "doc": f"/download/{rel}"}
        except Exception:
            return {"ok": True, "text": preview[:3000]}

    if "flashcards" in instruction or "cards" in instruction:
        cards = _make_flashcards(text, limit=10)
        preview = "\n\n".join([f"Q: {c['q']}\nA: {c['a']}" for c in cards])
        try:
            # Save as Q/A lines in DOCX
            lines = []
            for c in cards:
                lines.append(f"Q: {c['q']}")
                lines.append(f"A: {c['a']}")
                lines.append("")
            out = _generate_docx(f"Flashcards - {title}", lines)
            rel = os.path.relpath(out, os.getcwd()).replace("\\", "/")
            return {"ok": True, "text": preview[:3000], "cards": cards, "doc": f"/download/{rel}"}
        except Exception:
            return {"ok": True, "text": preview[:3000], "cards": cards}

    if "outline" in instruction:
        outline = _outline_text(text, max_sections=6)
        flat = []
        for sec in outline:
            flat.append(sec["title"]) 
            flat.extend(["  - " + b for b in sec["bullets"]])
            flat.append("")
        preview = "\n".join(flat)
        try:
            out = _generate_docx(f"Outline - {title}", flat)
            rel = os.path.relpath(out, os.getcwd()).replace("\\", "/")
            return {"ok": True, "text": preview[:3000], "doc": f"/download/{rel}"}
        except Exception:
            return {"ok": True, "text": preview[:3000]}

    if "clean" in instruction or "cleanup" in instruction:
        cleaned = _clean_text(text)
        try:
            out = _generate_docx(f"Cleaned - {title}", [cleaned])
            rel = os.path.relpath(out, os.getcwd()).replace("\\", "/")
            return {"ok": True, "text": cleaned[:3000], "doc": f"/download/{rel}"}
        except Exception:
            return {"ok": True, "text": cleaned[:3000]}

    if instruction.startswith("rewrite:") or instruction.startswith("tone:"):
        try:
            tone = instruction.split(":",1)[1].strip().lower()
        except Exception:
            tone = ""
        rewritten = _rewrite_tone(text, tone)
        try:
            out = _generate_docx(f"Rewritten ({tone or 'neutral'}) - {title}", [rewritten])
            rel = os.path.relpath(out, os.getcwd()).replace("\\", "/")
            return {"ok": True, "text": rewritten[:3000], "doc": f"/download/{rel}"}
        except Exception:
            return {"ok": True, "text": rewritten[:3000]}

    if "glossary" in instruction:
        items = _extract_glossary(text, limit=20)
        preview = "\n".join(items)
        try:
            out = _generate_docx(f"Glossary - {title}", items)
            rel = os.path.relpath(out, os.getcwd()).replace("\\", "/")
            return {"ok": True, "text": preview[:3000], "doc": f"/download/{rel}"}
        except Exception:
            return {"ok": True, "text": preview[:3000]}

    # default: return content and optionally repackage to docx
    try:
        out = _generate_docx(f"Processed - {title}", [text[:6000]])
        rel = os.path.relpath(out, os.getcwd()).replace("\\", "/")
        return {"ok": True, "text": text[:3000], "doc": f"/download/{rel}"}
    except Exception:
        return {"ok": True, "text": text[:3000]}



def _extract_text_auto(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()
    if ext == ".docx":
        return _extract_text_from_docx(path)
    if ext == ".pptx":
        return _extract_text_from_pptx(path)
    if ext == ".pdf":
        return _extract_text_from_pdf(path)
    return ""

def _sentences_set(text: str):
    sents = _split_sentences(text)
    # normalize
    return set([" ".join(s.lower().split()) for s in sents if s.strip()])

@app.route("/api/compare", methods=["POST"])
def api_compare():
    try:
        fA = request.files.get("fileA")
        fB = request.files.get("fileB")
        out_fmt = (request.form.get("format") or "both").lower()
        if not fA or not fB:
            return {"ok": False, "error": "need_two_files"}, 400
        nameA = secure_filename(fA.filename or "A")
        nameB = secure_filename(fB.filename or "B")
        pA = os.path.join(UPLOADS_DIR, f"cmp_A_{int(time.time())}_{nameA}")
        pB = os.path.join(UPLOADS_DIR, f"cmp_B_{int(time.time())}_{nameB}")
        fA.save(pA); fB.save(pB)
        tA = _extract_text_auto(pA)
        tB = _extract_text_auto(pB)
        if not tA and not tB:
            return {"ok": False, "error": "no_text_in_files"}, 400
        setA = _sentences_set(tA)
        setB = _sentences_set(tB)
        onlyA = [s for s in setA - setB][:20]
        onlyB = [s for s in setB - setA][:20]
        common = [s for s in setA & setB][:20]
        # Build preview and lines
        title = f"Compare - {os.path.splitext(nameA)[0]} vs {os.path.splitext(nameB)[0]}"
        lines = []
        lines.append("Summary:")
        lines.append(f"Only in {nameA}: {len(onlyA)}")
        lines.append(f"Only in {nameB}: {len(onlyB)}")
        lines.append(f"Common: {len(common)}")
        lines.append("")
        lines.append(f"Only in {nameA}:")
        lines += [" - " + s for s in onlyA]
        lines.append("")
        lines.append(f"Only in {nameB}:")
        lines += [" - " + s for s in onlyB]
        lines.append("")
        lines.append("Common Points:")
        lines += [" - " + s for s in common]
        lines.append("")
        # unified notes as key points from both texts
        uni = _extract_key_points((tA + "\n\n" + tB)[:16000], limit=12)
        lines.append("Unified Notes:")
        lines += [" - " + s for s in uni]
        preview = "\n".join(lines[:60])
        out = {"ok": True, "preview": preview[:3000]}
        if out_fmt in ("docx", "both"):
            try:
                doc_path = _generate_docx(title, lines)
                rel = os.path.relpath(doc_path, os.getcwd()).replace("\\", "/")
                out["doc"] = f"/download/{rel}"
            except Exception:
                pass
        if out_fmt in ("pdf", "both"):
            try:
                pdf_path = _generate_pdf(title, lines)
                rel = os.path.relpath(pdf_path, os.getcwd()).replace("\\", "/")
                out["pdf"] = f"/download/{rel}"
            except Exception:
                pass
        return out
    except Exception as e:
        return {"ok": False, "error": str(e)}, 500

@app.route("/api/merge", methods=["POST"])
def api_merge():
    try:
        out_fmt = (request.form.get("format") or "both").lower()
        # accept files as file1,file2,... or multiple under 'files'
        files = []
        if "files" in request.files:
            files = request.files.getlist("files")
        else:
            for i in range(1, 8):
                f = request.files.get(f"file{i}")
                if f: files.append(f)
        if len(files) < 2:
            return {"ok": False, "error": "need_two_or_more_files"}, 400
        merged_lines = []
        refs = []
        for idx, f in enumerate(files, 1):
            nm = secure_filename(f.filename or f"file{idx}")
            p = os.path.join(UPLOADS_DIR, f"merge_{int(time.time())}_{idx}_{nm}")
            f.save(p)
            txt = _extract_text_auto(p)
            refs.append(nm)
            if txt:
                merged_lines.append(f"=== {nm} ===")
                pts = _extract_key_points(txt, limit=12)
                merged_lines += [" - " + s for s in pts]
                merged_lines.append("")
        if not merged_lines:
            return {"ok": False, "error": "no_text_in_files"}, 400
        # De-dup lines
        seen = set(); uniq = []
        for ln in merged_lines:
            norm = " ".join(ln.lower().split())
            if norm in seen: continue
            seen.add(norm); uniq.append(ln)
        uniq.append(""); uniq.append("References:")
        uniq += [" - " + r for r in refs]
        title = "Merged Report"
        out = {"ok": True, "preview": "\n".join(uniq[:80])[:3000]}
        if out_fmt in ("docx", "both"):
            try:
                doc_path = _generate_docx(title, uniq)
                rel = os.path.relpath(doc_path, os.getcwd()).replace("\\", "/")
                out["doc"] = f"/download/{rel}"
            except Exception:
                pass
        if out_fmt in ("pdf", "both"):
            try:
                pdf_path = _generate_pdf(title, uniq)
                rel = os.path.relpath(pdf_path, os.getcwd()).replace("\\", "/")
                out["pdf"] = f"/download/{rel}"
            except Exception:
                pass
        return out
    except Exception as e:
        return {"ok": False, "error": str(e)}, 500
@app.route("/api/echo", methods=["POST"])
def api_echo():
    data = request.json or {}
    return {"ok": True, "echo": data}

@app.route("/api/command", methods=["POST"])
def api_command():
    data = request.json or {}
    text = str(data.get("text", ""))
    try:
        resp = process_command(text)
        # persist bot reply in user's chat history
        user_email = get_logged_in_email()
        if user_email and not isinstance(resp, dict):
            try:
                save_chat_line(user_email, "bot", str(resp))
            except Exception:
                pass
        if isinstance(resp, dict):
            out = {"ok": True}
            out.update(resp)
            return out
        return {"ok": True, "text": resp}
    except Exception as e:
        return {"ok": False, "error": str(e)}, 500

@socketio.on("connect")
def ws_connect():
    print("Client connected")
    emit("draco_response", {"text": "Draco is online and ready!"})

@socketio.on("disconnect")
def ws_disconnect():
    print("Client disconnected")

@socketio.on("user_command")
def ws_user_command(payload):
    try:
        text = payload.get("text", "")
        print("Received command from web:", text)
        response = process_command(text)
        # send a structured response (allow dict for web actions)
        if isinstance(response, dict):
            emit("draco_response", response)
        else:
            emit("draco_response", {"text": response})
        # persist bot reply for logged-in users
        user_email = get_logged_in_email()
        if user_email and not isinstance(response, dict):
            try:
                save_chat_line(user_email, "bot", str(response))
            except Exception:
                pass
    except Exception as e:
        print("Error handling user_command:", e)
        emit("draco_response", {"text": "Error processing command."})

# ------------- Background: optional voice listener thread -------------
if not ON_RENDER:
    # simple wake-word loop (optional)
    speak("Voice listener active.")
    while True:
        try:
            with sr.Microphone() as source:
                r.adjust_for_ambient_noise(source, duration=0.6)
                print("Listening for wake word...")
                audio = r.listen(source, timeout=4, phrase_time_limit=4)
                text = ""
                try:
                    text = r.recognize_google(audio, language="en-IN").lower()
                except Exception:
                    continue
                if "draco" in text:
                    speak("Yes? Say your command.")
                    audio2 = r.listen(source, timeout=6, phrase_time_limit=8)
                    try:
                        cmd = r.recognize_google(audio2, language="en-IN")
                        print("Heard command:", cmd)
                        res = process_command(cmd)
                        speak(res)
                    except Exception:
                        speak("I couldn't understand. Try again.") 
        except Exception as e:
            print("Voice loop error:", e)
            time.sleep(1)
else:
    print("Voice listener disabled on Render (no audio hardware).")
# ------------- Start-up -------------
if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))
    run_kwargs = {}
    if ON_SERVER:
        run_kwargs["allow_unsafe_werkzeug"] = True
    socketio.run(app, host="0.0.0.0", port=port, **run_kwargs)